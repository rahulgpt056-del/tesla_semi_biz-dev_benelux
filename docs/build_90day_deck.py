"""
Generate the Tesla Semi Benelux "First 90 Days" deck (.pptx).

Run:    python docs/build_90day_deck.py
Output: docs/Tesla-Semi-Benelux-First-90-Days.pptx
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(ROOT, "Tesla-Semi-Benelux-First-90-Days.pptx")

NAVY = RGBColor(0x0B, 0x25, 0x45)
RED = RGBColor(0xE3, 0x19, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x8A, 0x97, 0xA4)
INK = RGBColor(0x13, 0x31, 0x5C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = "Calibri"
    return box


def accent_bar(slide, left, top, width=0.12, height=1.2, color=RED):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bullets(slide, left, top, width, height, items, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


# --------------------------------------------------------------------------- #
# Slide 1 — Title
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, NAVY)
accent_bar(s, 0.7, 2.6, width=0.15, height=1.8)
add_text(s, 1.0, 2.5, 11, 1.2, "First 90 Days", 48, WHITE, bold=True)
add_text(s, 1.0, 3.5, 11, 0.8, "Tesla Semi · Benelux Business Development", 22, RED, bold=True)
add_text(
    s, 1.0, 4.4, 11, 1.2,
    "A practical plan for the Business Development Manager, Semi — Benelux role "
    "(Req. ID 267880), built around three working tools developed for this application.",
    14, GREY,
)

# --------------------------------------------------------------------------- #
# Slide 2 — The opportunity
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, WHITE)
accent_bar(s, 0.7, 0.6, height=0.9)
add_text(s, 1.0, 0.55, 11, 1.0, "The opportunity in Benelux", 30, NAVY, bold=True)
bullets(s, 1.0, 1.8, 11.3, 5.0, [
    "Regional, depot-return logistics networks (parcel, grocery, FMCG distribution) are a strong "
    "structural fit for Semi's ~800km range and overnight depot charging.",
    "The Netherlands' AanZET subsidy and the live vrachtwagenheffing (truck toll, with an expected "
    "zero-emission exemption) both push the diesel-vs-Semi economics in Semi's favour right now.",
    "Belgium's largest retail and logistics groups (e.g. Colruyt, H.Essers) are already publicly "
    "piloting alternative-fuel heavy trucks — the market is primed, not cold.",
    "The hard part isn't the vehicle decision — it's depot grid capacity and DSO connection lead "
    "times (Liander/Enexis/Stedin in NL, Fluvius/ORES in BE), which can take 6-18 months.",
    "This means the BD motion has to start the grid-readiness conversation in week one, not after "
    "a Sales Contract is signed.",
], size=16)

# --------------------------------------------------------------------------- #
# Slide 3 — Days 0-30
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, WHITE)
accent_bar(s, 0.7, 0.6, height=0.9, color=RED)
add_text(s, 1.0, 0.55, 11, 1.0, "Days 0-30 — Listen, qualify, build pipeline", 28, NAVY, bold=True)
bullets(s, 1.0, 1.8, 11.3, 5.0, [
    "Walk the Target-Account Territory Map (35 Benelux fleet operators, tiered A/B/C on duty-cycle "
    "fit) with the existing Benelux sales team to validate and reprioritise against live "
    "relationships and pipeline.",
    "Run first discovery calls with Tier A accounts — national parcel/grocery networks and groups "
    "with public sustainability mandates — focused on one representative route per fleet.",
    "For every qualifying conversation, plug real or estimated route data into the TCO & "
    "Duty-Cycle Calculator to get an immediate Fit / Conditional / No-fit verdict and an "
    "energy/charging snapshot.",
    "Open the grid-capacity conversation early: ask each account what their current depot "
    "connection capacity is, and whether they've had any contact with their DSO about EV charging.",
    "Deliverable by day 30: a qualified shortlist of 8-10 accounts with a duty-cycle fit verdict "
    "and an initial grid-capacity flag for each.",
], size=15)

# --------------------------------------------------------------------------- #
# Slide 4 — Days 31-60
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, WHITE)
accent_bar(s, 0.7, 0.6, height=0.9, color=RED)
add_text(s, 1.0, 0.55, 11, 1.0, "Days 31-60 — Spec the solution, get utility readiness moving", 26, NAVY, bold=True)
bullets(s, 1.0, 1.8, 11.3, 5.0, [
    "For each shortlisted account, use the calculator's Site & Charging view to size the fleet "
    "pilot: number of trucks, chargers needed, peak load vs. usable depot capacity.",
    "Where peak load exceeds usable capacity, help the account submit a DSO capacity-increase "
    "request immediately — this is the longest lead-time item and becomes the master schedule's "
    "critical path.",
    "Work the Site-Readiness Playbook stage by stage with each account's facilities/operations "
    "contact, in parallel with the commercial conversation (not sequentially after it).",
    "Begin assembling the incentive stack per account: AanZET subsidy status (current RVO budget "
    "round), expected toll treatment, and any regional charging-infrastructure co-funding.",
    "Deliverable by day 60: a sized vehicle + charging proposal and an open DSO connection request "
    "for at least 3-4 accounts.",
], size=15)

# --------------------------------------------------------------------------- #
# Slide 5 — Days 61-90
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, WHITE)
accent_bar(s, 0.7, 0.6, height=0.9, color=RED)
add_text(s, 1.0, 0.55, 11, 1.0, "Days 61-90 — Contract, onboarding, commissioning groundwork", 26, NAVY, bold=True)
bullets(s, 1.0, 1.8, 11.3, 5.0, [
    "Move the strongest 1-2 accounts toward a Sales Contract, using the TCO calculator's 5-year "
    "saving and payback figures as the core of the commercial narrative.",
    "Lock in incentive eligibility (AanZET budget round, toll exemption status) before the budget "
    "round risk becomes live — flagged explicitly in the playbook's risk register.",
    "Start onboarding planning with the account's operations team: driver familiarisation, depot "
    "charging schedule design, maintenance SLAs, telematics integration.",
    "Where a DSO connection request was filed in days 31-60, get a firm timeline and confirm "
    "whether a phased rollout (chargers within existing headroom) can start before the full "
    "upgrade completes.",
    "Deliverable by day 90: at least one account in contract or late-stage negotiation, with a "
    "site-readiness timeline that has utility, permitting, and onboarding workstreams already "
    "running in parallel — and a repeatable playbook for the next account on the territory map.",
], size=15)

# --------------------------------------------------------------------------- #
# Slide 6 — Tools built for this role
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, NAVY)
add_text(s, 1.0, 0.55, 11, 1.0, "Three tools, built and working", 28, WHITE, bold=True)
add_text(s, 1.0, 1.3, 11.3, 0.6,
         "Developed as part of this application — all open-source, unit-tested, and self-updating.",
         14, GREY)

cards = [
    ("TCO & Duty-Cycle Calculator", "Streamlit app — duty-cycle fit, energy & charging plan, "
     "5-year Semi-vs-diesel TCO and payback for any fleet. 22 unit tests."),
    ("Target-Account Territory Map", "Streamlit app — 35 Benelux fleet operators scored and "
     "tiered A/B/C on duty-cycle fit signals. 8 unit tests."),
    ("Site-Readiness Playbook", "This document's companion — a six-stage checklist from "
     "qualification through commissioning, with a risk register."),
]
left = 0.9
for title, desc in cards:
    box = s.shapes.add_shape(1, Inches(left), Inches(2.3), Inches(3.7), Inches(4.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x13, 0x31, 0x5C)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(13)
    r2.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE3)
    left += 4.0

# --------------------------------------------------------------------------- #
# Slide 7 — 90-day success metrics
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, WHITE)
accent_bar(s, 0.7, 0.6, height=0.9, color=RED)
add_text(s, 1.0, 0.55, 11, 1.0, "What success looks like at day 90", 28, NAVY, bold=True)
bullets(s, 1.0, 1.8, 11.3, 5.0, [
    "8-10 qualified accounts with a documented duty-cycle fit verdict and grid-capacity status.",
    "3-4 accounts with an open DSO connection request — the critical-path item started early.",
    "At least one account in contract or late-stage negotiation, with incentives locked in.",
    "A repeatable, tool-supported process the wider Benelux team can run on the next tier of "
    "accounts without rebuilding the analysis from scratch each time.",
], size=17)

# --------------------------------------------------------------------------- #
# Slide 8 — Closing
# --------------------------------------------------------------------------- #
s = add_slide()
fill_bg(s, NAVY)
accent_bar(s, 0.7, 2.8, width=0.15, height=1.2)
add_text(s, 1.0, 2.8, 11, 0.8, "Ready to start day one.", 32, WHITE, bold=True)
add_text(s, 1.0, 3.7, 11, 1.0,
         "Tesla Semi · Benelux  |  Business Development Manager, Semi — Benelux (Req. ID 267880)",
         14, GREY)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("WROTE", OUT)
